from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Color scheme
C_BG = RGBColor(0x0D, 0x1B, 0x2A)        # dark navy
C_ACCENT = RGBColor(0x00, 0xB4, 0xD8)     # cyan
C_ACCENT2 = RGBColor(0xFF, 0xA7, 0x00)    # amber
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xCA, 0xE9, 0xFF)
C_CARD = RGBColor(0x1A, 0x2F, 0x45)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def set_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, l, t, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txb

def make_title_slide(prs, week_num, title_th, title_en, clos, method, media, course="07-034-266 การสร้างธุรกิจเริ่มต้นด้วยนวัตกรรมและเทคโนโลยี"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    W, H = SLIDE_W, SLIDE_H

    # top accent bar
    add_rect(slide, 0, 0, W, Inches(0.07), C_ACCENT)

    # week badge circle (use rect rounded via shape)
    badge = add_rect(slide, Inches(0.5), Inches(0.3), Inches(1.8), Inches(1.8), C_ACCENT)
    add_text(slide, f"Week\n{week_num}", Inches(0.5), Inches(0.3), Inches(1.8), Inches(1.8), 22, C_BG, bold=True, align=PP_ALIGN.CENTER)

    # course name
    add_text(slide, course, Inches(2.5), Inches(0.35), Inches(10.3), Inches(0.5), 11, C_LIGHT, align=PP_ALIGN.LEFT)

    # divider line
    add_rect(slide, Inches(0.5), Inches(2.35), Inches(6), Inches(0.04), C_ACCENT)

    # Thai title
    add_text(slide, title_th, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.1), 32, C_WHITE, bold=True)

    # English title
    add_text(slide, title_en, Inches(0.5), Inches(3.55), Inches(12.3), Inches(0.6), 20, C_ACCENT)

    # CLOs card
    clo_box = add_rect(slide, Inches(0.5), Inches(4.3), Inches(5.8), Inches(2.7), C_CARD)
    add_text(slide, "CLOs ที่เกี่ยวข้อง", Inches(0.6), Inches(4.35), Inches(5.6), Inches(0.4), 13, C_ACCENT2, bold=True)
    add_text(slide, clos, Inches(0.6), Inches(4.75), Inches(5.6), Inches(2.1), 13, C_LIGHT)

    # Method card
    met_box = add_rect(slide, Inches(6.5), Inches(4.3), Inches(2.9), Inches(2.7), C_CARD)
    add_text(slide, "วิธีการสอน", Inches(6.6), Inches(4.35), Inches(2.7), Inches(0.4), 13, C_ACCENT2, bold=True)
    add_text(slide, method, Inches(6.6), Inches(4.75), Inches(2.7), Inches(2.1), 13, C_LIGHT)

    # Media card
    med_box = add_rect(slide, Inches(9.6), Inches(4.3), Inches(3.2), Inches(2.7), C_CARD)
    add_text(slide, "สื่อที่ใช้", Inches(9.7), Inches(4.35), Inches(3.0), Inches(0.4), 13, C_ACCENT2, bold=True)
    add_text(slide, media, Inches(9.7), Inches(4.75), Inches(3.0), Inches(2.1), 13, C_LIGHT)

    # bottom bar
    add_rect(slide, 0, H - Inches(0.4), W, Inches(0.4), C_CARD)
    add_text(slide, "อาจารย์ ดร.สุรเชษฐ์ สังขพันธ์  |  คณะวิทยาการจัดการ  |  มหาวิทยาลัยนราธิวาสราชนครินทร์",
             Inches(0.3), H - Inches(0.4), W - Inches(0.6), Inches(0.4), 10, C_LIGHT, align=PP_ALIGN.CENTER)

def add_content_slide(prs, title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    W, H = SLIDE_W, SLIDE_H

    add_rect(slide, 0, 0, W, Inches(0.07), C_ACCENT)
    add_rect(slide, 0, Inches(0.07), Inches(0.15), Inches(1.4), C_ACCENT)

    add_text(slide, title, Inches(0.4), Inches(0.15), W - Inches(0.7), Inches(1.3), 28, C_WHITE, bold=True)
    add_rect(slide, Inches(0.4), Inches(1.5), Inches(6), Inches(0.04), C_ACCENT)

    # bullet content area
    content_area = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), W - Inches(1.0), H - Inches(2.5))
    tf = content_area.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        if bullet.startswith("##"):
            run.text = bullet[2:].strip()
            run.font.size = Pt(18)
            run.font.color.rgb = C_ACCENT2
            run.font.bold = True
            p.space_before = Pt(14)
        elif bullet.startswith("#"):
            run.text = "▸  " + bullet[1:].strip()
            run.font.size = Pt(16)
            run.font.color.rgb = C_ACCENT
            run.font.bold = True
        else:
            run.text = "     •  " + bullet
            run.font.size = Pt(14)
            run.font.color.rgb = C_LIGHT

    if note:
        add_rect(slide, 0, H - Inches(0.8), W, Inches(0.04), C_ACCENT)
        add_text(slide, "💡 " + note, Inches(0.5), H - Inches(0.75), W - Inches(1.0), Inches(0.7), 12, C_ACCENT2)

    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), C_CARD)
    add_text(slide, "อาจารย์ ดร.สุรเชษฐ์ สังขพันธ์  |  คณะวิทยาการจัดการ  |  มหาวิทยาลัยนราธิวาสราชนครินทร์",
             Inches(0.3), H - Inches(0.35), W - Inches(0.6), Inches(0.35), 10, C_LIGHT, align=PP_ALIGN.CENTER)

def add_activity_slide(prs, week_num, activities):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    W, H = SLIDE_W, SLIDE_H

    add_rect(slide, 0, 0, W, Inches(0.07), C_ACCENT)
    add_text(slide, f"กิจกรรมสัปดาห์ที่ {week_num}", Inches(0.5), Inches(0.2), W - Inches(1.0), Inches(0.9), 28, C_WHITE, bold=True)
    add_rect(slide, Inches(0.4), Inches(1.1), Inches(4), Inches(0.04), C_ACCENT)

    card_colors = [C_ACCENT, C_ACCENT2, RGBColor(0x06, 0xD6, 0xA0), RGBColor(0xEF, 0x47, 0x6F)]
    n = len(activities)
    card_w = (W - Inches(1.0)) / n
    for i, (act_title, act_body) in enumerate(activities):
        x = Inches(0.5) + i * card_w
        add_rect(slide, x, Inches(1.3), card_w - Inches(0.15), Inches(5.7), C_CARD)
        add_rect(slide, x, Inches(1.3), card_w - Inches(0.15), Inches(0.08), card_colors[i % len(card_colors)])
        add_text(slide, act_title, x + Inches(0.1), Inches(1.45), card_w - Inches(0.35), Inches(0.6),
                 15, card_colors[i % len(card_colors)], bold=True)
        add_text(slide, act_body, x + Inches(0.1), Inches(2.1), card_w - Inches(0.35), Inches(4.7),
                 13, C_LIGHT)

    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), C_CARD)
    add_text(slide, "อาจารย์ ดร.สุรเชษฐ์ สังขพันธ์  |  คณะวิทยาการจัดการ  |  มหาวิทยาลัยนราธิวาสราชนครินทร์",
             Inches(0.3), H - Inches(0.35), W - Inches(0.6), Inches(0.35), 10, C_LIGHT, align=PP_ALIGN.CENTER)

def add_summary_slide(prs, week_num, takeaways, next_week, final=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    W, H = SLIDE_W, SLIDE_H

    add_rect(slide, 0, 0, W, Inches(0.07), C_ACCENT)
    add_text(slide, "สรุปและการบ้าน" if not final else "สรุปรายวิชา", Inches(0.5), Inches(0.2), Inches(8), Inches(0.9), 28, C_WHITE, bold=True)
    add_rect(slide, Inches(0.4), Inches(1.1), Inches(3.5), Inches(0.04), C_ACCENT)

    # Takeaways box
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(6.3), Inches(5.7), C_CARD)
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(6.3), Inches(0.08), C_ACCENT)
    add_text(slide, "Key Takeaways สิ่งที่ได้เรียนรู้", Inches(0.6), Inches(1.35), Inches(6.1), Inches(0.55), 16, C_ACCENT, bold=True)
    ta_text = "\n".join([f"✓  {t}" for t in takeaways])
    add_text(slide, ta_text, Inches(0.6), Inches(1.95), Inches(6.1), Inches(4.8), 14, C_LIGHT)

    # Next week box (or closing message on the final week)
    add_rect(slide, Inches(7.0), Inches(1.25), Inches(5.8), Inches(2.6), C_CARD)
    add_rect(slide, Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.08), C_ACCENT2)
    box1_label = "ก้าวต่อไป" if final else "สัปดาห์หน้า"
    add_text(slide, f"{box1_label}: {next_week}", Inches(7.1), Inches(1.35), Inches(5.6), Inches(2.4), 15, C_ACCENT2, bold=True)

    # Homework box (or course wrap-up on the final week)
    add_rect(slide, Inches(7.0), Inches(4.05), Inches(5.8), Inches(2.9), C_CARD)
    add_rect(slide, Inches(7.0), Inches(4.05), Inches(5.8), Inches(0.08), RGBColor(0x06, 0xD6, 0xA0))
    if final:
        add_text(slide, "สิ่งที่ต้องส่ง", Inches(7.1), Inches(4.15), Inches(5.6), Inches(0.5), 16, RGBColor(0x06, 0xD6, 0xA0), bold=True)
        add_text(slide, "ส่ง Personal Action Plan\nส่งแบบประเมินรายวิชาและอาจารย์\nขอบคุณที่ร่วมเดินทางตลอดภาคเรียน 🎓",
                 Inches(7.1), Inches(4.7), Inches(5.6), Inches(2.1), 13, C_LIGHT)
    else:
        add_text(slide, "การบ้าน / SDL", Inches(7.1), Inches(4.15), Inches(5.6), Inches(0.5), 16, RGBColor(0x06, 0xD6, 0xA0), bold=True)
        add_text(slide, f"ทบทวนเนื้อหาสัปดาห์นี้\nเตรียมพร้อมสำหรับ: {next_week}\nศึกษาด้วยตนเอง 5 ชั่วโมง",
                 Inches(7.1), Inches(4.7), Inches(5.6), Inches(2.1), 13, C_LIGHT)

    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), C_CARD)
    add_text(slide, "อาจารย์ ดร.สุรเชษฐ์ สังขพันธ์  |  คณะวิทยาการจัดการ  |  มหาวิทยาลัยนราธิวาสราชนครินทร์",
             Inches(0.3), H - Inches(0.35), W - Inches(0.6), Inches(0.35), 10, C_LIGHT, align=PP_ALIGN.CENTER)

# =====================================================================
# WEEK DATA
# =====================================================================

weeks = [
    # ---- WEEK 1 ----
    {
        "week": 1,
        "title_th": "แนะนำรายวิชาและแนวคิดการเป็นผู้ประกอบการ ธุรกิจเริ่มต้น (Startup)",
        "title_en": "Course Introduction & Entrepreneurship / Startup Concepts",
        "clos": "CLO1: แสดงพฤติกรรมซื่อสัตย์ มีวินัย ตรงต่อเวลา\nCLO2: อธิบายแนวคิดการเป็นผู้ประกอบการ",
        "method": "บรรยาย\nอภิปราย\nกรณีศึกษา",
        "media": "สไลด์\nกรณีศึกษา\nVideo Clip",
        "slides": [
            ("ภาพรวมรายวิชา", [
                "#รหัส 07-034-266 / 3 หน่วยกิต (2-2-5)",
                "วิสาหกิจเริ่มต้น (Startup) คืออะไร",
                "ทำไมต้องเรียนรายวิชานี้ในยุคดิจิทัล",
                "#กติกาการเรียนและการประเมิน",
                "การเข้าเรียน วินัย ตรงต่อเวลา",
                "งานกลุ่ม / แผนธุรกิจ / Pitch Deck",
                "เกณฑ์การให้คะแนน (100 คะแนน)",
            ]),
            ("ผู้ประกอบการคือใคร?", [
                "#นิยามผู้ประกอบการ (Entrepreneur)",
                "บุคคลที่สร้างคุณค่าใหม่ด้วยการรับความเสี่ยง",
                "มองเห็นโอกาสที่คนอื่นมองข้าม",
                "#ความแตกต่าง: ธุรกิจทั่วไป vs Startup",
                "Startup: เติบโตเร็ว, ขยายได้ (Scalable)",
                "ธุรกิจทั่วไป: เติบโตตามเส้นตรง",
                "#ตัวอย่าง Startup ไทยที่ประสบความสำเร็จ",
                "Grab, LINE MAN, Omise, Wongnai",
            ], "Startup ไม่ใช่แค่ App มือถือ — มันคือวิธีคิดแบบใหม่"),
            ("Startup Ecosystem ในประเทศไทย", [
                "#ระบบนิเวศ Startup ไทย",
                "NIA (สำนักงานนวัตกรรมแห่งชาติ)",
                "depa (สำนักงานส่งเสริมเศรษฐกิจดิจิทัล)",
                "สสว. (สำนักงานส่งเสริมวิสาหกิจขนาดกลางและขนาดย่อม)",
                "#แหล่งทุนและการสนับสนุน",
                "Angel Investor, Venture Capital (VC)",
                "โครงการ Startup Thailand",
                "#เทรนด์ Startup โลกปี 2025-2026",
                "AI-First Startup, Green Tech, HealthTech",
            ]),
        ],
        "activities": [
            ("บรรยาย", "อาจารย์แนะนำรายวิชา\nชี้แจงกติกา เกณฑ์การประเมิน\nและแผนการสอนตลอดภาคเรียน"),
            ("อภิปรายกลุ่ม", "นักศึกษาแบ่งกลุ่ม\nระดมความคิด:\n'Startup ที่รู้จัก คืออะไร\nทำไมถึงสำเร็จหรือล้มเหลว?'"),
            ("กรณีศึกษา", "ศึกษา Startup ไทย 1 ราย\nวิเคราะห์จุดเด่น\nและนำเสนอหน้าชั้นเรียน"),
        ],
        "takeaways": [
            "เข้าใจความหมายของ Startup และผู้ประกอบการ",
            "รู้จักระบบนิเวศ Startup ในไทยและเทรนด์โลก",
            "เข้าใจกติกาและเกณฑ์การประเมินรายวิชา",
            "เริ่มคิดถึงปัญหาในชีวิตประจำวันที่อยากแก้ไข",
        ],
        "next_week": "นวัตกรรมและความคิดสร้างสรรค์ทางธุรกิจ / Design Thinking",
    },
    # ---- WEEK 2 ----
    {
        "week": 2,
        "title_th": "นวัตกรรมและความคิดสร้างสรรค์ทางธุรกิจ กระบวนการคิดเชิงออกแบบ",
        "title_en": "Business Innovation & Creativity / Design Thinking Process",
        "clos": "CLO2: อธิบายนวัตกรรมทางธุรกิจ\nCLO3: ค้นหาและพัฒนาไอเดียนวัตกรรม",
        "method": "บรรยาย\nเวิร์กชอป\nกิจกรรมกลุ่ม",
        "media": "สไลด์\nใบงาน Design Thinking\nVideo",
        "slides": [
            ("นวัตกรรมคืออะไร?", [
                "#นิยามนวัตกรรม (Innovation)",
                "ไอเดียใหม่ที่สร้างคุณค่าและถูกนำไปใช้จริง",
                "ไม่ใช่แค่การประดิษฐ์ — ต้องมี Impact",
                "#ประเภทนวัตกรรม",
                "Product Innovation: สินค้า/บริการใหม่",
                "Process Innovation: กระบวนการใหม่",
                "Business Model Innovation: โมเดลธุรกิจใหม่",
                "#ตัวอย่าง: Netflix เปลี่ยน Business Model",
                "จาก DVD rental → Streaming subscription",
            ], "นวัตกรรมไม่จำเป็นต้องเป็นเทคโนโลยีเสมอไป"),
            ("Design Thinking 5 ขั้นตอน", [
                "#1. Empathize — เข้าใจผู้ใช้",
                "สัมภาษณ์, สังเกต, สวมบทบาท (จากมุมมองผู้ใช้)",
                "#2. Define — นิยามปัญหา",
                "Point of View (POV) Statement",
                "#3. Ideate — ระดมไอเดีย",
                "Brainstorming, SCAMPER, How Might We",
                "#4. Prototype — สร้างต้นแบบ",
                "ต้นแบบอย่างรวดเร็ว (Low-Fidelity)",
                "#5. Test — ทดสอบ",
                "รับ Feedback จากผู้ใช้จริง → ปรับปรุง",
            ]),
            ("เครื่องมือ Ideation", [
                "#Brainstorming Rules",
                "ไม่ตัดสินไอเดียในช่วงระดม",
                "ยิ่งมากยิ่งดี (Quantity over Quality)",
                "#SCAMPER Technique",
                "S=Substitute, C=Combine, A=Adapt",
                "M=Modify, P=Put to other use, E=Eliminate, R=Rearrange",
                "#How Might We (HMW) Questions",
                "เปลี่ยนปัญหาเป็นคำถามเปิด",
                "เช่น: 'เราจะช่วยให้คนในชนบทเข้าถึงอาหารสดได้อย่างไร?'",
            ], "Design Thinking เป็นกระบวนการวนซ้ำ ไม่ใช่เส้นตรง"),
        ],
        "activities": [
            ("บรรยาย", "อธิบาย Design Thinking\n5 ขั้นตอนพร้อมตัวอย่าง\nจาก Startup จริง\n(IDEO, Airbnb)"),
            ("เวิร์กชอป", "ฝึก Empathy Map:\nนักศึกษาเลือก 1 กลุ่มเป้าหมาย\nและวาด Empathy Map\nบน Whiteboard/ใบงาน"),
            ("ใบงาน", "กลุ่มส่ง POV Statement\n'[ผู้ใช้] ต้องการ [ความต้องการ]\nเพราะ [ข้อมูลเชิงลึก]'\nนำเสนอสั้นๆ 2 นาที"),
        ],
        "takeaways": [
            "เข้าใจความหมายและประเภทของนวัตกรรมทางธุรกิจ",
            "รู้จัก Design Thinking 5 ขั้นตอน",
            "ฝึกเขียน Empathy Map และ POV Statement",
            "เริ่มมองปัญหาจากมุมมองของผู้ใช้งาน",
        ],
        "next_week": "การค้นหาปัญหา/โอกาสทางธุรกิจ และการพัฒนาไอเดียนวัตกรรม (Ideation)",
    },
    # ---- WEEK 3 ----
    {
        "week": 3,
        "title_th": "การค้นหาปัญหา/โอกาสทางธุรกิจ และการพัฒนาไอเดียนวัตกรรม",
        "title_en": "Problem Discovery & Business Opportunity / Ideation",
        "clos": "CLO3: ค้นหาและพัฒนาไอเดียนวัตกรรม\nCLO5: ทำงานร่วมกับผู้อื่นในกลุ่ม",
        "method": "ปฏิบัติการกลุ่ม\nระดมสมอง\nนำเสนอ",
        "media": "ใบงาน\nเครื่องมือระดมสมอง\nกระดาน Post-it",
        "slides": [
            ("Pain Points และโอกาสทางธุรกิจ", [
                "#Pain Points คืออะไร?",
                "ปัญหาที่ผู้คนพบในชีวิตประจำวัน",
                "ความไม่สะดวก ความเสียเวลา ความสูญเสียเงิน",
                "#3 แหล่งหา Pain Points",
                "สังเกตชีวิตประจำวัน (Observation)",
                "สัมภาษณ์กลุ่มเป้าหมาย (Interview)",
                "ข้อมูล Online Reviews / Social Media",
                "#Pain Point → โอกาส Startup",
                "Grab: ปัญหาแท็กซี่ไม่พอ → แอปเรียกรถ",
                "Foodpanda: ปัญหาสั่งอาหาร → Delivery App",
            ], "ปัญหาที่ดีที่สุดคือปัญหาที่คุณเองก็เจอ"),
            ("Idea Generation Techniques", [
                "#Blue Ocean Strategy",
                "สร้างตลาดใหม่ที่ไม่มีคู่แข่ง",
                "ERRC Framework: Eliminate, Reduce, Raise, Create",
                "#Jobs-to-be-Done Theory",
                "ลูกค้าซื้อ 'งาน' ไม่ใช่ 'สินค้า'",
                "เช่น: ซื้อสว่านเพื่อได้รู ไม่ใช่ตัวสว่าน",
                "#Trend Analysis",
                "Megatrend: AI, Aging Society, Green Economy",
                "ใช้ Google Trends, TikTok Trends",
            ]),
            ("เกณฑ์ประเมินไอเดีย Startup", [
                "#Problem-Solution Fit",
                "ปัญหาชัดเจนและใหญ่พอ?",
                "Solution แก้ปัญหาได้จริง?",
                "#3 มิติประเมินไอเดีย",
                "Desirability: คนต้องการไหม?",
                "Feasibility: ทำได้จริงไหม?",
                "Viability: ทำแล้วมีกำไรไหม?",
                "#เครื่องมือ: Idea Evaluation Matrix",
                "ให้คะแนน 1-5 ในแต่ละมิติ",
                "เลือกไอเดียที่คะแนนรวมสูงสุด",
            ], "ไอเดียที่ดีไม่จำเป็นต้องเป็นไอเดียแรก — ลองหลายๆ ไอเดียก่อน"),
        ],
        "activities": [
            ("Problem Hunting", "นักศึกษาออกสำรวจ\nบริเวณมหาวิทยาลัย\nค้นหา 5 Pain Points\nบันทึกลงใบงาน"),
            ("Brainstorming", "กลุ่มระดมไอเดีย\nด้วย Post-it Note\nไม่ตัดสินไอเดีย\nนาน 15 นาที"),
            ("Idea Pitch", "แต่ละกลุ่มเลือก\nไอเดียดีที่สุด 1 ข้อ\nนำเสนอ 3 นาที:\nปัญหา / กลุ่มเป้าหมาย / Solution"),
        ],
        "takeaways": [
            "สามารถค้นหา Pain Points จากสภาพแวดล้อมรอบตัว",
            "รู้จักเทคนิค Idea Generation (Blue Ocean, JTBD)",
            "ประเมินไอเดียด้วย Desirability-Feasibility-Viability",
            "เริ่มมีไอเดีย Startup เบื้องต้นสำหรับโครงงาน",
        ],
        "next_week": "การวิเคราะห์ตลาดและความต้องการของลูกค้า (Customer & Market Validation)",
    },
    # ---- WEEK 4 ----
    {
        "week": 4,
        "title_th": "การวิเคราะห์ตลาดและความต้องการของลูกค้า",
        "title_en": "Customer & Market Validation",
        "clos": "CLO3: ค้นหาและพัฒนาไอเดียนวัตกรรม\nCLO5: ทำงานร่วมกับผู้อื่นในกลุ่ม",
        "method": "บรรยาย\nปฏิบัติ\nสำรวจตลาด",
        "media": "สไลด์\nแบบสำรวจ\nGoogle Forms",
        "slides": [
            ("ทำไมต้องทำ Market Validation?", [
                "#ปัญหาของ Startup ส่วนใหญ่",
                "42% ของ Startup ล้มเหลวเพราะ 'ไม่มีคนต้องการ'",
                "สร้างสิ่งที่ตัวเองอยากสร้าง ≠ สิ่งที่ตลาดต้องการ",
                "#Build-Measure-Learn Loop",
                "สร้าง → วัดผล → เรียนรู้ → ปรับปรุง",
                "Lean Startup Methodology (Eric Ries)",
                "#Customer Discovery Process",
                "กำหนดสมมติฐาน (Assumption)",
                "ออกพบลูกค้า (Get out of the building!)",
            ], "'ลูกค้าไม่เคยโกหก' — ฟังให้มากกว่าพูด"),
            ("Customer Persona", [
                "#Customer Persona คืออะไร?",
                "โปรไฟล์ตัวแทนลูกค้าในอุดมคติ",
                "#องค์ประกอบ Persona",
                "ชื่อ / อายุ / อาชีพ / รายได้",
                "ปัญหา / ความต้องการ / พฤติกรรม",
                "ช่องทางที่ใช้ (Social Media, Apps)",
                "#วิธีเก็บข้อมูลสร้าง Persona",
                "สัมภาษณ์เชิงลึก 5-10 คน",
                "Google Analytics / Social Insight",
                "แบบสอบถามออนไลน์",
            ]),
            ("TAM SAM SOM Framework", [
                "#TAM (Total Addressable Market)",
                "ขนาดตลาดทั้งหมดในโลก",
                "เช่น: ตลาด Food Delivery ไทย = 80,000 ล้านบาท/ปี",
                "#SAM (Serviceable Addressable Market)",
                "ส่วนของตลาดที่เราเข้าถึงได้",
                "เช่น: ภาคใต้ + กลุ่ม Gen Z = 5,000 ล้านบาท/ปี",
                "#SOM (Serviceable Obtainable Market)",
                "ส่วนตลาดที่จะ 'ได้จริง' ใน 1-3 ปีแรก",
                "เช่น: 1% ของ SAM = 50 ล้านบาท/ปี",
            ], "SOM ต้องสมเหตุสมผล — นักลงทุนดูตรงนี้เป็นพิเศษ"),
        ],
        "activities": [
            ("Mini Survey", "กลุ่มออกแบบ\nแบบสำรวจ 5-8 ข้อ\nใน Google Forms\nเกี่ยวกับไอเดียของกลุ่ม"),
            ("สัมภาษณ์", "สัมภาษณ์นักศึกษา\nและบุคลากร 5 คน\nบันทึกข้อมูลจริง\n(ไม่ใช่การสมมติ)"),
            ("วิเคราะห์", "สรุปผล Survey\nสร้าง Customer Persona\nและประมาณ TAM/SAM/SOM\nนำเสนอ 5 นาที"),
        ],
        "takeaways": [
            "เข้าใจว่า Market Validation สำคัญกว่าการสร้างสินค้าทันที",
            "สร้าง Customer Persona จากข้อมูลจริง",
            "ประมาณขนาดตลาดด้วย TAM/SAM/SOM",
            "ฝึกทักษะการสัมภาษณ์และการฟัง",
        ],
        "next_week": "แบบจำลองธุรกิจ Business Model Canvas (BMC)",
    },
    # ---- WEEK 5 ----
    {
        "week": 5,
        "title_th": "แบบจำลองธุรกิจ (Business Model Canvas)",
        "title_en": "Business Model Canvas (BMC)",
        "clos": "CLO3: จัดทำและนำเสนอแผนธุรกิจ\nCLO5: ทำงานร่วมกับผู้อื่นในกลุ่ม",
        "method": "เวิร์กชอป\nกิจกรรมกลุ่ม",
        "media": "BMC Template\nPost-it\nDigital Canvas",
        "slides": [
            ("Business Model Canvas คืออะไร?", [
                "#BMC: เครื่องมือออกแบบโมเดลธุรกิจ",
                "พัฒนาโดย Alexander Osterwalder",
                "1 หน้ากระดาษ แทน Business Plan 30 หน้า",
                "#9 Building Blocks",
                "1. Customer Segments (CS)",
                "2. Value Propositions (VP)",
                "3. Channels (CH)",
                "4. Customer Relationships (CR)",
                "5. Revenue Streams (RS)",
                "6. Key Resources (KR)",
                "7. Key Activities (KA)",
                "8. Key Partnerships (KP)",
                "9. Cost Structure (CS)",
            ], "BMC เป็นภาพรวมธุรกิจทั้งหมดในหน้าเดียว"),
            ("4 ด้านหลักของ BMC", [
                "##ด้านขวา: Customer Side (รายรับ)",
                "#Customer Segments",
                "กลุ่มลูกค้าที่เจาะจง — B2C, B2B, Niche",
                "#Value Propositions",
                "คุณค่าที่เราให้กับลูกค้า (ทำไมเลือกเรา?)",
                "#Channels & Customer Relationships",
                "ช่องทางถึงลูกค้า / วิธีรักษาความสัมพันธ์",
                "##ด้านซ้าย: Company Side (ต้นทุน)",
                "#Key Resources / Activities / Partnerships",
                "ทรัพยากร กิจกรรมหลัก และพันธมิตร",
                "#Cost Structure",
                "ต้นทุนคงที่และต้นทุนผันแปร",
            ]),
            ("Revenue Models ที่พบใน Startup", [
                "#Subscription Model",
                "Netflix, Spotify — จ่ายรายเดือน",
                "#Freemium Model",
                "LINE, Canva — ใช้ฟรี + Premium features",
                "#Marketplace Model",
                "Shopee, Grab — ค่าธรรมเนียมต่อธุรกรรม",
                "#SaaS (Software as a Service)",
                "B2B software — บริการซอฟต์แวร์ออนไลน์",
                "#Advertising Model",
                "Facebook, TikTok — รายได้จากโฆษณา",
            ], "เลือก Revenue Model ที่เหมาะกับธุรกิจ อย่าลอกแบบ"),
        ],
        "activities": [
            ("Workshop BMC", "แต่ละกลุ่มรับ\nBMC Template\n(กระดาษ A1 หรือ Miro)\nกรอก 9 ช่องให้ครบ"),
            ("Gallery Walk", "ติด BMC บนผนัง\nกลุ่มอื่น ๆ เดินดู\nและแปะ Post-it\nข้อดี / ข้อสงสัย"),
            ("Feedback", "แต่ละกลุ่มนำเสนอ BMC\n5 นาที + รับ Feedback\nจากอาจารย์และเพื่อน\nปรับปรุง BMC"),
        ],
        "takeaways": [
            "เข้าใจ 9 Building Blocks ของ Business Model Canvas",
            "สามารถออกแบบ BMC สำหรับ Startup ของกลุ่มได้",
            "รู้จัก Revenue Models แบบต่างๆ",
            "ฝึกให้ Feedback และรับ Feedback อย่างสร้างสรรค์",
        ],
        "next_week": "คุณค่าที่นำเสนอ (Value Proposition) และ MVP",
    },
    # ---- WEEK 6 ----
    {
        "week": 6,
        "title_th": "คุณค่าที่นำเสนอ (Value Proposition) และต้นแบบผลิตภัณฑ์ขั้นต่ำ (MVP)",
        "title_en": "Value Proposition & Minimum Viable Product (MVP)",
        "clos": "CLO3: จัดทำและนำเสนอแผนธุรกิจ\nCLO4: ผลิตสื่อและชิ้นงานโดยใช้เครื่องมือดิจิทัล",
        "method": "ปฏิบัติการกลุ่ม\nออกแบบต้นแบบ",
        "media": "ใบงาน\nเครื่องมือต้นแบบ (Figma/Canva)",
        "slides": [
            ("Value Proposition Canvas", [
                "#Value Proposition Canvas (VPC)",
                "เชื่อมโยง Customer Profile ↔ Value Map",
                "#Customer Profile (วงกลมขวา)",
                "Jobs: งานที่ลูกค้าต้องทำให้สำเร็จ",
                "Pains: ความเจ็บปวด/อุปสรรค",
                "Gains: ผลลัพธ์ที่ต้องการ",
                "#Value Map (สี่เหลี่ยมซ้าย)",
                "Products & Services: สิ่งที่เรานำเสนอ",
                "Pain Relievers: บรรเทาความเจ็บปวด",
                "Gain Creators: สร้างผลลัพธ์ที่ต้องการ",
            ], "Product-Market Fit เกิดเมื่อ Value Map ตรงกับ Customer Profile"),
            ("Unique Value Proposition (UVP)", [
                "#UVP คืออะไร?",
                "ประโยคเดียวที่บอกว่าเราช่วยใคร ทำอะไร ต่างจากใคร",
                "#สูตร UVP ง่ายๆ",
                "เราช่วย [กลุ่มเป้าหมาย]",
                "ให้ [คุณประโยชน์หลัก]",
                "ต่างจากคู่แข่งตรงที่ [ความแตกต่าง]",
                "#ตัวอย่าง UVP จาก Startup จริง",
                "Slack: 'Be more productive at work with less effort'",
                "Airbnb: 'Belong anywhere'",
                "Grab: 'Everyday Everything App'",
            ]),
            ("MVP: Minimum Viable Product", [
                "#MVP คืออะไร?",
                "ผลิตภัณฑ์ขั้นต่ำที่ทดสอบสมมติฐานหลักได้",
                "ใช้ทรัพยากรน้อยที่สุด เรียนรู้ให้เร็วที่สุด",
                "#ประเภท MVP",
                "Landing Page MVP (เว็บเพจง่ายๆ)",
                "Concierge MVP (ทำด้วยมือก่อน)",
                "Wizard of Oz MVP (จำลองระบบ)",
                "Paper Prototype (ต้นแบบกระดาษ)",
                "#กรณีศึกษา: Dropbox MVP",
                "แค่วิดีโอ 3 นาที → มีคนลงทะเบียน 75,000 คน",
            ], "MVP ≠ สินค้าที่ด้อยคุณภาพ — แต่คือการเรียนรู้ที่เร็วที่สุด"),
        ],
        "activities": [
            ("VPC Workshop", "กลุ่มกรอก\nValue Proposition Canvas\nสำหรับ Startup ของตน\nบน Miro หรือใบงาน"),
            ("UVP Writing", "แต่ละกลุ่มเขียน UVP\nในเวลา 10 นาที\nนำเสนอและรับ Feedback\nจากเพื่อนในชั้น"),
            ("MVP Design", "ออกแบบ MVP:\nเลือกประเภท MVP\nวาด Wireframe/Sketch\nหรือสร้าง Landing Page\nด้วย Canva"),
        ],
        "takeaways": [
            "เข้าใจ Value Proposition Canvas และหาจุดตรงกัน",
            "เขียน UVP ที่ชัดเจนสำหรับ Startup ของกลุ่ม",
            "รู้จักประเภท MVP และวิธีเลือกให้เหมาะสม",
            "เริ่มออกแบบ Prototype เบื้องต้นของ Startup",
        ],
        "next_week": "ระบบโลจิสติกส์ การขนส่ง และคลังสินค้าสำหรับธุรกิจเริ่มต้น",
    },
    # ---- WEEK 7 ----
    {
        "week": 7,
        "title_th": "ระบบโลจิสติกส์ การขนส่ง และคลังสินค้าสำหรับธุรกิจเริ่มต้น",
        "title_en": "Logistics, Transportation & Warehouse Management for Startups",
        "clos": "CLO2: อธิบายพื้นฐานระบบโลจิสติกส์ การขนส่ง และคลังสินค้า",
        "method": "บรรยาย\nกรณีศึกษา\nอภิปราย",
        "media": "สไลด์\nกรณีศึกษา\nVideo",
        "slides": [
            ("โลจิสติกส์และ Supply Chain", [
                "#โลจิสติกส์ (Logistics) คืออะไร?",
                "การวางแผนและจัดการการเคลื่อนย้ายสินค้า",
                "จากต้นทางไปยังปลายทางอย่างมีประสิทธิภาพ",
                "#Supply Chain Management (SCM)",
                "Supplier → ผู้ผลิต → คลังสินค้า → ขนส่ง → ลูกค้า",
                "#ความสำคัญต่อ Startup",
                "ลด Cost ได้มากถึง 30-50% จากโลจิสติกส์ที่ดี",
                "Customer Experience: ส่งเร็ว ส่งถูก ส่งตรง",
            ], "โลจิสติกส์ที่ดีคือ Competitive Advantage ที่มองไม่เห็น"),
            ("การจัดการคลังสินค้า", [
                "#Warehouse Management System (WMS)",
                "ระบบจัดการสินค้าในคลังสินค้าด้วยซอฟต์แวร์",
                "#หลักการ FIFO / FEFO",
                "FIFO: First In, First Out",
                "FEFO: First Expired, First Out (อาหาร/ยา)",
                "#Inventory Management",
                "Safety Stock: สต๊อกสำรองฉุกเฉิน",
                "Reorder Point: จุดที่ต้องสั่งสินค้าใหม่",
                "#Just-In-Time (JIT)",
                "สั่งสินค้าเมื่อต้องการ ลดต้นทุนการเก็บสต๊อก",
            ]),
            ("Last-Mile Delivery และ Startup", [
                "#Last-Mile Delivery",
                "การส่งสินค้าจากคลังถึงมือลูกค้า",
                "ต้นทุนสูงสุด 53% ของค่าขนส่งทั้งหมด",
                "#โมเดลส่งสินค้าสำหรับ Startup ไทย",
                "3PL (Third-Party Logistics): Flash, Kerry, J&T",
                "Dropshipping: ไม่ต้องมีสต๊อก",
                "Fulfillment Center: ฝากเก็บ/แพ็ค/ส่ง",
                "#E-commerce Logistics Trend",
                "Same-day delivery, Drone delivery",
                "Sustainable Packaging",
            ], "Startup ส่วนใหญ่ใช้ 3PL ก่อน — ไม่ต้องลงทุนคลังสินค้าเอง"),
        ],
        "activities": [
            ("บรรยาย", "อธิบายระบบโลจิสติกส์\nSupply Chain\nและการจัดการคลังสินค้า\nพร้อม Diagram"),
            ("กรณีศึกษา", "วิเคราะห์กรณีศึกษา:\n'Startup อาหารสด'\nมีปัญหาโลจิสติกส์อะไร?\nแก้ไขอย่างไร?"),
            ("อภิปรายกลุ่ม", "แต่ละกลุ่มวาง\nStrategy โลจิสติกส์\nสำหรับ Startup ของตน\n(ใช้ 3PL ตัวไหน? ทำไม?)"),
        ],
        "takeaways": [
            "เข้าใจระบบ Supply Chain และโลจิสติกส์พื้นฐาน",
            "รู้จักการจัดการคลังสินค้าและ Inventory",
            "เลือกโมเดลขนส่งที่เหมาะกับ Startup ของตนได้",
            "เข้าใจความสำคัญของ Last-Mile Delivery ต่อ Customer Experience",
        ],
        "next_week": "การวางแผนการเงินและการประเมินความเป็นไปได้ของธุรกิจ",
    },
    # ---- WEEK 8 ----
    {
        "week": 8,
        "title_th": "การวางแผนการเงินและการประเมินความเป็นไปได้ของธุรกิจ",
        "title_en": "Financial Planning & Business Feasibility Assessment",
        "clos": "CLO3: จัดทำและนำเสนอแผนธุรกิจและกลยุทธ์การจัดการ",
        "method": "บรรยาย\nปฏิบัติ\nวิเคราะห์ตัวเลข",
        "media": "สไลด์\nตารางการเงิน (Excel/Sheets)",
        "slides": [
            ("Financial Planning สำหรับ Startup", [
                "#ทำไม Startup ต้องวางแผนการเงิน?",
                "38% ของ Startup ล้มเหลวเพราะเงินหมด (Cash Flow)",
                "นักลงทุนดูตัวเลขการเงินก่อนตัดสินใจ",
                "#งบการเงินหลัก 3 รายการ",
                "Income Statement: รายได้ - ค่าใช้จ่าย = กำไร",
                "Balance Sheet: สินทรัพย์ = หนี้สิน + ทุน",
                "Cash Flow Statement: เงินสดเข้า-ออก",
                "#Unit Economics",
                "CAC (Customer Acquisition Cost): ต้นทุนหาลูกค้าใหม่ 1 ราย",
                "LTV (Lifetime Value): มูลค่าตลอดชีวิตของลูกค้า",
            ], "LTV ต้องมากกว่า CAC อย่างน้อย 3 เท่า — กฎทอง Startup"),
            ("การประมาณการรายได้และค่าใช้จ่าย", [
                "#Revenue Projection (3 ปี)",
                "Year 1: ทดสอบตลาด — ไม่ต้องกำไร",
                "Year 2: เติบโต — เริ่มคุ้มทุน",
                "Year 3: Scale — กำไรชัดเจน",
                "#ต้นทุนหลักของ Startup",
                "Fixed Cost: เงินเดือน ค่าเช่า ค่าซอฟต์แวร์",
                "Variable Cost: ต้นทุนสินค้า ค่าขนส่ง",
                "#Break-Even Analysis",
                "จุดคุ้มทุน = Fixed Cost ÷ (ราคา - VC ต่อหน่วย)",
                "ต้องขายเท่าไหร่จึงจะคุ้มทุน?",
            ]),
            ("Funding และแหล่งเงินทุน Startup", [
                "#Bootstrapping",
                "ใช้เงินตัวเอง / เพื่อน / ครอบครัว (FFF)",
                "ควบคุมได้ 100% แต่เงินจำกัด",
                "#Grant และทุนภาครัฐ",
                "NIA: ทุน Startup ไทย (สูงสุด 1 ล้านบาท)",
                "depa: ทุนดิจิทัล SME",
                "#Angel Investor",
                "บุคคลที่ลงทุนในระยะแรก",
                "ได้เงิน + Mentor + เครือข่าย",
                "#Venture Capital (VC)",
                "ระยะ Series A ขึ้นไป — Valuation สูง",
            ], "เริ่มด้วย Bootstrapping หรือ Grant — อย่ารีบขายหุ้น"),
        ],
        "activities": [
            ("Workshop Excel", "กลุ่มสร้าง\nFinancial Model\nบน Google Sheets:\nรายได้ / ค่าใช้จ่าย / กำไร\n3 ปีข้างหน้า"),
            ("Break-Even", "คำนวณจุดคุ้มทุน\nของ Startup กลุ่มตนเอง\nต้องขายกี่หน่วย?\nใช้เวลากี่เดือน?"),
            ("Pitch Finance", "นำเสนอตัวเลขการเงิน\n5 สไลด์:\nRevenue / Cost / Profit\nBreak-Even / Funding Plan"),
        ],
        "takeaways": [
            "เข้าใจงบการเงินหลัก 3 รายการสำหรับ Startup",
            "คำนวณ CAC, LTV และ Break-Even Point ได้",
            "วาง Revenue Projection 3 ปีให้สมเหตุสมผล",
            "รู้จักแหล่งทุนและเลือกแหล่งทุนที่เหมาะกับระยะของธุรกิจ",
        ],
        "next_week": "กลยุทธ์การตลาดดิจิทัลและการเติบโตของธุรกิจ (Growth)",
    },
    # ---- WEEK 9 ----
    {
        "week": 9,
        "title_th": "กลยุทธ์การตลาดดิจิทัลและการเติบโตของธุรกิจ (Growth)",
        "title_en": "Digital Marketing Strategy & Growth",
        "clos": "CLO3: จัดทำและนำเสนอแผนธุรกิจและกลยุทธ์การจัดการ",
        "method": "บรรยาย\nอภิปราย\nกรณีศึกษา",
        "media": "สไลด์\nกรณีศึกษา",
        "slides": [
            ("Digital Marketing สำหรับ Startup", [
                "#ทำไม Digital Marketing สำคัญกับ Startup",
                "ต้นทุนต่ำกว่าการตลาดแบบดั้งเดิม วัดผลได้ทันที",
                "เข้าถึงกลุ่มเป้าหมายเฉพาะกลุ่มได้แม่นยำ",
                "#ช่องทางหลัก",
                "Social Media Marketing (Facebook, IG, TikTok)",
                "Search Engine Marketing (SEO/SEM)",
                "Email Marketing & Influencer Marketing",
                "#กรณีศึกษา: Wongnai",
                "สร้าง Community รีวิวร้านอาหาร ก่อนขยายสู่ Ads/Delivery",
            ], "เลือกช่องทางที่ลูกค้าเป้าหมายอยู่ ไม่ใช่ทุกช่องทาง"),
            ("Brand Identity และ Content Marketing", [
                "#Brand Identity คืออะไร",
                "ตัวตนของแบรนด์: โลโก้ สี Tone of Voice",
                "ความสอดคล้องสร้างความน่าเชื่อถือ",
                "#Content Marketing Pillars",
                "Educate / Entertain / Inspire / Convert",
                "#กรณีศึกษา: Instagram",
                "สร้าง Brand ผ่าน Visual-first Content ตั้งแต่วันแรก",
                "#เครื่องมือวัดผล Content",
                "Engagement Rate, Reach, Conversion Rate",
            ]),
            ("AARRR Growth Funnel (Pirate Metrics)", [
                "#AARRR Framework (Dave McClure)",
                "Acquisition: ลูกค้ามาจากไหน",
                "Activation: ลูกค้าได้ประสบการณ์ที่ดีครั้งแรกไหม",
                "Retention: ลูกค้ากลับมาใช้ซ้ำไหม",
                "Revenue: ลูกค้าจ่ายเงินไหม",
                "Referral: ลูกค้าบอกต่อไหม",
                "#กรณีศึกษา: Uber Referral Program",
                "ให้ Credit ทั้งผู้แนะนำและผู้ถูกแนะนำ → Growth Loop",
            ], "Growth Hacking คือการทดลองอย่างเป็นระบบ ไม่ใช่เวทมนตร์"),
        ],
        "activities": [
            ("วิเคราะห์ Brand", "เลือก Startup 1 ราย\nวิเคราะห์ Brand Identity\nและช่องทาง Digital Marketing\nที่ใช้"),
            ("ออกแบบ Content Plan", "กลุ่มออกแบบ\nContent Calendar 1 สัปดาห์\nสำหรับ Startup ของตน\n(3 โพสต์ + เป้าหมาย)"),
            ("AARRR Mapping", "กลุ่มระบุกลยุทธ์\nในแต่ละขั้นของ AARRR\nสำหรับ Startup ตนเอง\nนำเสนอ 3 นาที"),
        ],
        "takeaways": [
            "เข้าใจช่องทาง Digital Marketing ที่เหมาะกับ Startup",
            "ออกแบบ Brand Identity และ Content Marketing เบื้องต้นได้",
            "ใช้ AARRR Framework วางกลยุทธ์การเติบโต",
            "เห็นตัวอย่างจริงจาก Wongnai, Instagram, Uber",
        ],
        "next_week": "การใช้ปัญญาประดิษฐ์และเครื่องมือดิจิทัลในการสร้างสื่อและพัฒนาธุรกิจ",
    },
    # ---- WEEK 10 ----
    {
        "week": 10,
        "title_th": "การใช้ปัญญาประดิษฐ์และเครื่องมือดิจิทัลในการสร้างสื่อและพัฒนาธุรกิจ",
        "title_en": "AI & Digital Tools for Content Creation and Business Development",
        "clos": "CLO4: ผลิตสื่อและชิ้นงานนำเสนอแผนธุรกิจโดยใช้เครื่องมือดิจิทัลและปัญญาประดิษฐ์",
        "method": "บรรยาย\nเวิร์กชอป",
        "media": "สไลด์\nใบงาน",
        "slides": [
            ("AI สำหรับ Startup ยุคใหม่", [
                "#ทำไม Startup ต้องใช้ AI",
                "ลดต้นทุน เพิ่มความเร็วในการทำงาน แข่งกับธุรกิจใหญ่ได้",
                "#ประเภทเครื่องมือ AI ที่ใช้บ่อย",
                "Generative AI: ChatGPT, Claude, Gemini",
                "Design AI: Canva AI, Midjourney",
                "Video/Voice AI: CapCut AI, ElevenLabs",
            ], "AI คือผู้ช่วย ไม่ใช่ผู้ตัดสินใจแทนผู้ประกอบการ"),
            ("ใช้ AI สร้างสื่อธุรกิจ", [
                "#Content Creation ด้วย AI",
                "เขียน Caption, Script, Blog ด้วย Prompt ที่ดี",
                "ออกแบบโลโก้/โปสเตอร์ด้วย Canva AI",
                "#หลักการเขียน Prompt ที่ดี",
                "ระบุ Context, Format, Tone, ตัวอย่าง",
                "#ข้อควรระวัง",
                "ตรวจสอบความถูกต้อง (Hallucination)",
                "ลิขสิทธิ์และความเป็นต้นฉบับ",
            ]),
            ("AI ในการพัฒนาธุรกิจ", [
                "#วิเคราะห์ข้อมูลด้วย AI",
                "สรุปผล Survey, วิเคราะห์ Feedback ลูกค้า",
                "#AI Chatbot สำหรับลูกค้า",
                "ตอบคำถามอัตโนมัติ ลดต้นทุน Customer Service",
                "#เครื่องมือ No-Code/AI สร้าง MVP เร็วขึ้น",
                "Bubble, Lovable, Replit Agent",
            ], "ใช้ AI เร่งความเร็ว แต่ Core Value ต้องมาจากความเข้าใจลูกค้าจริง"),
        ],
        "activities": [
            ("Prompt Workshop", "ฝึกเขียน Prompt\nสร้าง Caption โฆษณา\nสำหรับ Startup ของกลุ่ม\nด้วย ChatGPT/Claude"),
            ("AI Design", "ใช้ Canva AI\nออกแบบโลโก้/โปสเตอร์\nสำหรับแบรนด์ของกลุ่ม"),
            ("นำเสนอผลงาน", "นำเสนอสื่อที่สร้าง\nด้วย AI 1 ชิ้น\nพร้อมอธิบาย Prompt\nที่ใช้"),
        ],
        "takeaways": [
            "รู้จักเครื่องมือ AI หลักสำหรับสร้างสื่อและพัฒนาธุรกิจ",
            "เขียน Prompt ที่มีประสิทธิภาพได้",
            "ใช้ AI สร้างสื่อนำเสนอ/โลโก้/Content ได้จริง",
            "เข้าใจข้อจำกัดและข้อควรระวังในการใช้ AI",
        ],
        "next_week": "การจัดทำแผนธุรกิจ (Business Plan)",
    },
    # ---- WEEK 11 ----
    {
        "week": 11,
        "title_th": "การจัดทำแผนธุรกิจ (Business Plan)",
        "title_en": "Business Plan Development",
        "clos": "CLO3: จัดทำและนำเสนอแผนธุรกิจและกลยุทธ์การจัดการ\nCLO4: ผลิตสื่อและชิ้นงานนำเสนอแผนธุรกิจ",
        "method": "ปฏิบัติการกลุ่ม\nระดมสมอง",
        "media": "ใบงาน\nเครื่องมือระดมสมอง",
        "slides": [
            ("โครงสร้างแผนธุรกิจ (Business Plan)", [
                "#องค์ประกอบหลักของ Business Plan",
                "Executive Summary, Problem/Solution, Market Analysis",
                "Business Model, Marketing Plan, Financial Plan, Team",
                "#ความยาวที่เหมาะสม",
                "10-20 หน้า ไม่ใช่ 100 หน้า — นักลงทุนอ่านไม่ทัน",
            ], "Business Plan ที่ดีคือเอกสารที่ 'มีชีวิต' ปรับปรุงได้เสมอ"),
            ("Startup Funding Stages — เส้นทางเงินทุน", [
                "#Pre-Seed",
                "ไอเดียและทีม ใช้เงินส่วนตัว/FFF ทดสอบสมมติฐาน",
                "#Seed",
                "มี MVP/Early Traction ระดมจาก Angel/Seed VC",
                "#Series A",
                "Product-Market Fit ชัดเจน ขยายทีมและตลาด",
                "#Series B/C และต่อไป",
                "Scale ธุรกิจ ขยายต่างประเทศ/ควบรวมกิจการ",
                "#สิ่งที่นักลงทุนแต่ละระยะมองหา",
                "Pre-Seed: ทีมและไอเดีย / Series A+: ตัวเลขและการเติบโต",
            ], "อย่าระดมทุนเกินความจำเป็นของระยะธุรกิจตนเอง (Don't over-raise)"),
            ("เขียน Business Plan ให้น่าสนใจ", [
                "#เล่าเรื่องด้วย Storytelling",
                "เริ่มจาก Pain Point ที่คนเข้าใจง่าย",
                "#ใส่ตัวเลขที่ตรวจสอบได้",
                "อ้างอิงแหล่งข้อมูลจริง ไม่ประมาณการลอยๆ",
                "#กรณีศึกษา: Airbnb",
                "เริ่มจาก Business Plan เล็กๆ เช่าที่นอนลม สู่ Unicorn",
            ], "นักลงทุนซื้อ 'ทีมที่แก้ปัญหาได้' มากกว่า 'ไอเดียที่สมบูรณ์แบบ'"),
        ],
        "activities": [
            ("ระดมสมอง", "กลุ่มร่าง Outline\nBusiness Plan 8 หัวข้อ\nโดยใช้ข้อมูลจาก\nสัปดาห์ที่ผ่านมา"),
            ("กำหนด Funding Stage", "วิเคราะห์ว่า Startup\nของกลุ่มอยู่ระยะใด\n(Pre-Seed/Seed)\nและควรระดมทุนเท่าไหร่"),
            ("Peer Review", "แลกเปลี่ยน Outline\nกับกลุ่มอื่น\nให้ Feedback\nเพื่อปรับปรุงก่อนสัปดาห์ Pitch"),
        ],
        "takeaways": [
            "เข้าใจโครงสร้างและองค์ประกอบของ Business Plan",
            "รู้จักเส้นทางเงินทุน Startup ตั้งแต่ Pre-Seed ถึง Series C",
            "ระบุระยะการระดมทุนที่เหมาะกับ Startup ของกลุ่มได้",
            "ร่าง Business Plan Outline ของกลุ่มตนเองได้",
        ],
        "next_week": "การพัฒนาสื่อนำเสนอ (Pitch Deck) และเทคนิคการนำเสนอแผนธุรกิจ",
    },
    # ---- WEEK 12 ----
    {
        "week": 12,
        "title_th": "การพัฒนาสื่อนำเสนอ (Pitch Deck) และเทคนิคการนำเสนอแผนธุรกิจ",
        "title_en": "Pitch Deck Development & Investor Presentation Techniques",
        "clos": "CLO4: ผลิตสื่อและชิ้นงานนำเสนอแผนธุรกิจโดยใช้เครื่องมือดิจิทัลและปัญญาประดิษฐ์\nCLO5: ทำงานร่วมกับผู้อื่นในกลุ่ม",
        "method": "บรรยาย\nปฏิบัติ",
        "media": "สไลด์\nแบบสำรวจ",
        "slides": [
            ("โครงสร้าง Pitch Deck มาตรฐาน", [
                "#10-12 สไลด์ที่นักลงทุนอยากเห็น",
                "Cover, Problem, Solution, Market Size",
                "Business Model, Traction, Competition",
                "Team, Financials, Funding Ask, Vision",
                "#กรณีศึกษา: Airbnb Pitch Deck (2008)",
                "เรียบง่าย ใช้ข้อมูลจริง สื่อสารตรงประเด็น",
            ], "Pitch Deck ที่ดี อ่านได้ใน 3 นาทีโดยไม่ต้องมีคนอธิบาย"),
            ("Investor Relations — สื่อสารกับนักลงทุน", [
                "#นักลงทุนประเภทต่างๆ",
                "Angel Investor, VC, Government Grant, Crowdfunding",
                "#สิ่งที่นักลงทุนถามบ่อย",
                "'ทำไมต้องเป็นทีมคุณ?' 'Traction ตอนนี้เป็นอย่างไร?'",
                "#Cap Table และ Equity เบื้องต้น",
                "สัดส่วนการถือหุ้นเปลี่ยนแปลงตามรอบระดมทุน",
                "#หลัง Pitch: Due Diligence และ Term Sheet",
                "ขั้นตอนก่อนปิดดีลจริง",
            ], "ความสัมพันธ์กับนักลงทุนไม่จบที่ Pitch — ต้อง Update สม่ำเสมอ"),
            ("เทคนิคการนำเสนอ (Presentation Skills)", [
                "#Storytelling Structure",
                "Hook → Problem → Solution → Proof → Ask",
                "#Delivery Techniques",
                "น้ำเสียง ภาษากาย การสบตา จับเวลา",
                "#รับมือคำถาม Q&A",
                "ตอบตรงประเด็น ยอมรับถ้าไม่รู้ พร้อม Follow-up",
                "#ฝึกซ้อมก่อนนำเสนอจริง",
                "Practice Pitch กับเพื่อนกลุ่มอื่น",
            ], "การฝึกซ้อมพูด 10 ครั้ง สำคัญกว่าสไลด์สวย 1 ชุด"),
        ],
        "activities": [
            ("สร้าง Pitch Deck", "กลุ่มสร้าง Pitch Deck\n10-12 สไลด์\nด้วย Canva/PowerPoint\nตามโครงสร้างมาตรฐาน"),
            ("ฝึกพูด Pitch", "ฝึกพูด Pitch\nภายในกลุ่ม จับเวลา 3 นาที\nรับ Feedback จากเพื่อน"),
            ("Mock Q&A", "จำลองสถานการณ์\nนักลงทุนถามคำถาม\nฝึกตอบกระชับ\nและมั่นใจ"),
        ],
        "takeaways": [
            "สร้าง Pitch Deck ตามโครงสร้างมาตรฐานสากลได้",
            "เข้าใจกระบวนการ Investor Relations และ Cap Table เบื้องต้น",
            "ฝึกเทคนิคการนำเสนอและการตอบคำถาม Q&A",
            "พร้อมสำหรับการนำเสนอโครงงานจริงในสัปดาห์ถัดไป",
        ],
        "next_week": "ปฏิบัติการกลุ่ม: พัฒนาโครงงานธุรกิจเริ่มต้น และจรรยาบรรณในการทำงาน",
    },
    # ---- WEEK 13 ----
    {
        "week": 13,
        "title_th": "ปฏิบัติการกลุ่ม: พัฒนาโครงงานธุรกิจเริ่มต้น และจรรยาบรรณในการทำงาน",
        "title_en": "Group Workshop: Startup Project Development & Professional Ethics",
        "clos": "CLO1: แสดงพฤติกรรมซื่อสัตย์ มีวินัย ตรงต่อเวลา\nCLO5: ทำงานร่วมกับผู้อื่นในกลุ่ม",
        "method": "เวิร์กชอป",
        "media": "BMC Template",
        "slides": [
            ("Workshop: ปรับปรุงโครงงานธุรกิจ", [
                "#ทบทวนและขัดเกลา BMC, แผนธุรกิจ, Pitch Deck",
                "รวมข้อมูล Feedback จากสัปดาห์ที่ผ่านมา",
                "#จุดที่กลุ่มมักพลาด",
                "ตัวเลขการเงินไม่สมเหตุสมผล / ไม่มี Traction",
                "Value Proposition ไม่ชัดเจน",
            ]),
            ("การจดทะเบียนธุรกิจและกฎระเบียบที่เกี่ยวข้อง", [
                "#รูปแบบธุรกิจที่จดทะเบียนได้",
                "บุคคลธรรมดา / ห้างหุ้นส่วน / บริษัทจำกัด",
                "#หน่วยงานสนับสนุน Startup ภาครัฐ",
                "NIA, depa, สสว., BOI (สิทธิประโยชน์ทางภาษี)",
                "#กฎระเบียบที่ต้องรู้",
                "PDPA (พ.ร.บ.คุ้มครองข้อมูลส่วนบุคคล)",
                "ทรัพย์สินทางปัญญา (เครื่องหมายการค้า/สิทธิบัตร)",
            ], "การจดทะเบียนถูกต้องตั้งแต่ต้น ป้องกันปัญหาใหญ่ในอนาคต"),
            ("จรรยาบรรณผู้ประกอบการ (Business Ethics)", [
                "#ความซื่อสัตย์ต่อลูกค้าและนักลงทุน",
                "ไม่กล่าวเกินจริงเรื่อง Traction/ตัวเลข",
                "#ความรับผิดชอบต่อสังคมและสิ่งแวดล้อม (ESG)",
                "#การทำงานเป็นทีมอย่างมีวินัย",
                "แบ่งงานเป็นธรรม ตรงต่อเวลา ให้เครดิตเพื่อนร่วมทีม",
            ], "ชื่อเสียงของผู้ประกอบการสร้างยากแต่เสียง่าย — ซื่อสัตย์เสมอ"),
        ],
        "activities": [
            ("Workshop ปรับปรุง", "กลุ่มปรับปรุง\nBMC/แผนธุรกิจ/Pitch Deck\nตาม Feedback ที่ได้รับ"),
            ("ตรวจสอบกฎระเบียบ", "ตรวจสอบว่า Startup\nของกลุ่มต้องจดทะเบียน\nหรือขออนุญาตอะไรบ้าง"),
            ("ซ้อมใหญ่ (Dress Rehearsal)", "ซ้อม Pitch แบบเต็มรูปแบบ\nจับเวลาจริง\nก่อนนำเสนอสัปดาห์หน้า"),
        ],
        "takeaways": [
            "ปรับปรุงโครงงานธุรกิจให้พร้อมนำเสนอจริง",
            "รู้จักรูปแบบการจดทะเบียนธุรกิจและหน่วยงานสนับสนุนภาครัฐ",
            "เข้าใจกฎระเบียบพื้นฐาน เช่น PDPA และทรัพย์สินทางปัญญา",
            "ตระหนักถึงจรรยาบรรณผู้ประกอบการในการทำงานจริง",
        ],
        "next_week": "นำเสนอโครงงานธุรกิจ (Pitching) และวิจารณ์/ให้ข้อเสนอแนะ",
    },
    # ---- WEEK 14 ----
    {
        "week": 14,
        "title_th": "นำเสนอโครงงานธุรกิจ (Pitching) และวิจารณ์/ให้ข้อเสนอแนะ",
        "title_en": "Final Project Pitching & Peer Critique",
        "clos": "CLO4: ผลิตสื่อและชิ้นงานนำเสนอแผนธุรกิจโดยใช้เครื่องมือดิจิทัลและปัญญาประดิษฐ์\nCLO5: ทำงานร่วมกับผู้อื่นในกลุ่ม",
        "method": "ปฏิบัติการ",
        "media": "ใบงาน\nเครื่องมือต้นแบบ",
        "slides": [
            ("วันนำเสนอโครงงานธุรกิจจริง", [
                "#รูปแบบการนำเสนอ",
                "แต่ละกลุ่ม Pitch 7-10 นาที + Q&A 5 นาที",
                "#เกณฑ์การประเมิน (Rubric)",
                "Problem-Solution Fit, Business Model, การนำเสนอ, ความคิดสร้างสรรค์",
            ], "วันนี้คือผลรวมของการเรียนรู้ตลอดภาคเรียน — นำเสนอด้วยความมั่นใจ"),
            ("บทเรียนจาก Startup ที่สำเร็จและล้มเหลว", [
                "#PayPal Mafia",
                "ทีมผู้ก่อตั้งที่แตกตัวไปสร้าง Tesla, LinkedIn, YouTube",
                "#Uber: Growth ที่รวดเร็วแต่มีข้อโต้เถียงด้านกฎระเบียบ",
                "#บทเรียนจากความล้มเหลว",
                "ธุรกิจที่ไม่ฟังตลาด หรือหมดเงินก่อนถึง Product-Market Fit",
            ], "ทุก Pitch วันนี้คือจุดเริ่มต้น ไม่ใช่จุดสิ้นสุดของไอเดีย"),
            ("การให้และรับ Feedback อย่างสร้างสรรค์", [
                "#หลักการให้ Feedback (Sandwich Method)",
                "ชื่นชม → ข้อเสนอแนะ → ให้กำลังใจ",
                "#การรับ Feedback อย่างมีวุฒิภาวะ",
                "ฟังโดยไม่ตั้งรับ บันทึกประเด็นสำคัญ",
                "#นำ Feedback ไปพัฒนาต่อ",
                "ปรับปรุง Business Plan ฉบับสมบูรณ์ส่งท้ายภาค",
            ]),
        ],
        "activities": [
            ("Pitching จริง", "แต่ละกลุ่มนำเสนอ\nPitch Deck 7-10 นาที\nต่อหน้าอาจารย์และเพื่อน"),
            ("Q&A Session", "รับคำถามจาก\nอาจารย์และกลุ่มอื่น\nฝึกตอบอย่างมั่นใจ"),
            ("ให้ Feedback เพื่อน", "แต่ละคนให้ Feedback\nกลุ่มอื่นอย่างน้อย 1 กลุ่ม\nตามแบบฟอร์มที่กำหนด"),
        ],
        "takeaways": [
            "นำเสนอโครงงานธุรกิจฉบับสมบูรณ์ต่อหน้าผู้ฟังจริง",
            "ฝึกตอบคำถามและรับมือสถานการณ์เฉพาะหน้า",
            "เรียนรู้จากกรณีศึกษาความสำเร็จและความล้มเหลวของ Startup จริง",
            "ให้และรับ Feedback อย่างสร้างสรรค์",
        ],
        "next_week": "สะท้อนผลการเรียนรู้ ตั้งเป้าหมายพัฒนาตนเอง และสรุปรายวิชา",
    },
    # ---- WEEK 15 ----
    {
        "week": 15,
        "title_th": "สะท้อนผลการเรียนรู้ ตั้งเป้าหมายพัฒนาตนเอง และสรุปรายวิชา",
        "title_en": "Learning Reflection, Personal Goal Setting & Course Wrap-up",
        "clos": "CLO6: สะท้อนผลการเรียนรู้และพัฒนาตนเองอย่างต่อเนื่อง",
        "method": "อภิปราย\nสะท้อนผล",
        "media": "แบบสะท้อนการเรียนรู้",
        "slides": [
            ("ทบทวนเส้นทางการเรียนรู้ตลอดภาคเรียน", [
                "#จาก Week 1 ถึง Week 14",
                "จากไอเดียเริ่มต้น สู่แผนธุรกิจและ Pitch Deck ฉบับสมบูรณ์",
                "#ทักษะที่ได้พัฒนา",
                "Design Thinking, BMC, Financial Planning, Digital Marketing, Pitching",
            ], "การเป็นผู้ประกอบการคือการเรียนรู้ตลอดชีวิต ไม่ใช่จบที่วิชานี้"),
            ("เส้นทางต่อจากนี้สำหรับไอเดียของคุณ", [
                "#ถ้าอยากทำต่อจริง",
                "สมัครโครงการบ่มเพาะ (Incubator/Accelerator)",
                "ทุน NIA, depa, มหาวิทยาลัย",
                "#ถ้ายังไม่พร้อมทำ Startup ตอนนี้",
                "นำทักษะผู้ประกอบการไปใช้ในงานประจำ (Intrapreneurship)",
                "#Community และเครือข่าย",
                "Startup Thailand, Techsauce, กลุ่ม Founder ท้องถิ่น",
            ]),
            ("ตั้งเป้าหมายพัฒนาตนเอง (Personal Action Plan)", [
                "#สะท้อนตนเอง 3 คำถาม",
                "ได้เรียนรู้อะไรเกี่ยวกับตัวเองบ้าง?",
                "จุดแข็ง-จุดที่ต้องพัฒนาในการเป็นผู้ประกอบการ?",
                "ขั้นตอนต่อไปใน 6 เดือนข้างหน้าคืออะไร?",
                "#เขียน Personal Action Plan",
                "เป้าหมาย / ขั้นตอน / Timeline / ผู้สนับสนุน",
            ], "เริ่มจากก้าวเล็กๆ ที่ทำได้จริงสัปดาห์นี้"),
        ],
        "activities": [
            ("สะท้อนผลกลุ่ม", "อภิปรายกลุ่มย่อย:\nบทเรียนสำคัญที่สุด\nที่ได้จากวิชานี้"),
            ("เขียน Action Plan", "แต่ละคนเขียน\nPersonal Action Plan\nส่งในแบบสะท้อนการเรียนรู้"),
            ("ประเมินรายวิชา", "ตอบแบบประเมิน\nรายวิชาและอาจารย์\nเพื่อพัฒนาการสอนต่อไป"),
        ],
        "takeaways": [
            "ทบทวนและเชื่อมโยงความรู้ตลอดภาคเรียนเป็นภาพรวม",
            "รู้จักเส้นทางต่อยอด เช่น Incubator, ทุนภาครัฐ, Community",
            "เขียน Personal Action Plan เพื่อพัฒนาตนเองต่อ",
            "สามารถสรุปและประเมินคุณค่าที่ได้รับจากรายวิชา",
        ],
        "next_week": "ขอให้ทุกคนนำความรู้ไปต่อยอดและประสบความสำเร็จในการเป็นผู้ประกอบการ!",
        "final": True,
    },
]

# =====================================================================
# GENERATE FILES
# =====================================================================
out_dir = "/Users/surachetsungkhapan/Startup_07034266/slides"
os.makedirs(out_dir, exist_ok=True)

for w in weeks:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # blank layout
    while len(prs.slide_layouts) < 7:
        pass  # already has layouts

    # Slide 1: Title
    make_title_slide(prs, w["week"], w["title_th"], w["title_en"], w["clos"], w["method"], w["media"])

    # Content slides
    for (title, bullets, *rest) in w["slides"]:
        note = rest[0] if rest else ""
        add_content_slide(prs, title, bullets, note)

    # Activity slide
    add_activity_slide(prs, w["week"], w["activities"])

    # Summary slide
    add_summary_slide(prs, w["week"], w["takeaways"], w["next_week"], final=w.get("final", False))

    filename = os.path.join(out_dir, f"Week{w['week']:02d}_Startup.pptx")
    prs.save(filename)
    print(f"Saved: {filename}")

print("\nDone! All 15 slides created.")
